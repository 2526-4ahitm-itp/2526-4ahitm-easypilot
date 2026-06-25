#!/usr/bin/env python3
"""Generates branded icon tiles and builds the EasyPilot iOS PowerPoint deck."""
import os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
LOGOS = os.path.join(HERE, "logos")
os.makedirs(LOGOS, exist_ok=True)

# ---- Theme (matches the EasyPilot app) ----
BG     = RGBColor(0x0A, 0x0F, 0x1E)
CARD   = RGBColor(0x14, 0x1C, 0x33)
ACCENT = RGBColor(0x29, 0x79, 0xFF)
TEXT   = RGBColor(0xE9, 0xEE, 0xFB)
MUTED  = RGBColor(0x8A, 0x97, 0xB8)
GREEN  = RGBColor(0x33, 0xD9, 0x80)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
CODEBG = RGBColor(0x0C, 0x13, 0x26)

SF = "/System/Library/Fonts/SFNS.ttf"
MONO = "/System/Library/Fonts/SFNSMono.ttf"

def font(path, size):
    return ImageFont.truetype(path, size)

# ---------------------------------------------------------------------------
# Icon tile generator: rounded gradient square + simple glyph + label
# ---------------------------------------------------------------------------
def rounded_tile(name, c1, c2, draw_glyph):
    S = 512
    img = Image.new("RGBA", (S, S), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    # vertical gradient background on a rounded square
    grad = Image.new("RGBA", (S, S))
    gd = ImageDraw.Draw(grad)
    for y in range(S):
        t = y / S
        r = int(c1[0] + (c2[0] - c1[0]) * t)
        g = int(c1[1] + (c2[1] - c1[1]) * t)
        b = int(c1[2] + (c2[2] - c1[2]) * t)
        gd.line([(0, y), (S, y)], fill=(r, g, b, 255))
    mask = Image.new("L", (S, S), 0)
    ImageDraw.Draw(mask).rounded_rectangle([0, 0, S - 1, S - 1], radius=110, fill=255)
    img.paste(grad, (0, 0), mask)
    d = ImageDraw.Draw(img)
    draw_glyph(d, S)
    img.save(os.path.join(LOGOS, f"{name}.png"))

def g_swiftui(d, S):
    # paint swirl suggestion: white rounded "play/leaf" + UI text
    f = font(SF, 150)
    d.text((S/2, S*0.42), "{ }", font=font(MONO, 150), fill=(255,255,255,255), anchor="mm")
    d.text((S/2, S*0.74), "SwiftUI", font=font(SF, 70), fill=(255,255,255,235), anchor="mm")

def g_scenekit(d, S):
    # isometric cube
    cx, cy, r = S/2, S*0.40, 120
    top = [(cx, cy-r), (cx+r, cy-r/2), (cx, cy), (cx-r, cy-r/2)]
    left = [(cx-r, cy-r/2), (cx, cy), (cx, cy+r), (cx-r, cy+r/2)]
    right = [(cx+r, cy-r/2), (cx, cy), (cx, cy+r), (cx+r, cy+r/2)]
    d.polygon(top,  fill=(255,255,255,255))
    d.polygon(left, fill=(255,255,255,150))
    d.polygon(right,fill=(255,255,255,200))
    d.text((S/2, S*0.80), "SceneKit", font=font(SF, 64), fill=(255,255,255,235), anchor="mm")

def g_coremotion(d, S):
    cx, cy = S/2, S*0.40
    for rr, w in [(150,14),(110,12),(70,10)]:
        d.ellipse([cx-rr, cy-rr*0.55, cx+rr, cy+rr*0.55], outline=(255,255,255,255), width=w)
    d.ellipse([cx-18, cy-18, cx+18, cy+18], fill=(255,255,255,255))
    d.text((S/2, S*0.80), "CoreMotion", font=font(SF, 56), fill=(255,255,255,235), anchor="mm")

def g_network(d, S):
    cx, cy = S/2, S*0.52
    for i, rr in enumerate([60, 120, 180]):
        bb = [cx-rr, cy-rr, cx+rr, cy+rr]
        d.arc(bb, start=215, end=325, fill=(255,255,255,255), width=16)
    d.ellipse([cx-16, cy-16, cx+16, cy+16], fill=(255,255,255,255))
    d.text((S/2, S*0.82), "Network", font=font(SF, 60), fill=(255,255,255,235), anchor="mm")

def g_websocket(d, S):
    # two opposing arrows
    cy = S*0.40
    d.line([(S*0.20, cy), (S*0.80, cy)], fill=(255,255,255,255), width=18)
    d.polygon([(S*0.80,cy),(S*0.68,cy-34),(S*0.68,cy+34)], fill=(255,255,255,255))
    cy2 = S*0.56
    d.line([(S*0.20, cy2), (S*0.80, cy2)], fill=(255,255,255,255), width=18)
    d.polygon([(S*0.20,cy2),(S*0.32,cy2-34),(S*0.32,cy2+34)], fill=(255,255,255,255))
    d.text((S/2, S*0.82), "WebSocket", font=font(SF, 56), fill=(255,255,255,235), anchor="mm")

rounded_tile("swiftui",    (0x0A,0x84,0xFF),(0x00,0x4F,0xC2), g_swiftui)
rounded_tile("scenekit",   (0x8E,0x5C,0xFF),(0x52,0x2C,0xB0), g_scenekit)
rounded_tile("coremotion", (0x33,0xD9,0x80),(0x12,0x8A,0x4E), g_coremotion)
rounded_tile("network",    (0x29,0x79,0xFF),(0x14,0x46,0xB0), g_network)
rounded_tile("websocket",  (0xFF,0x99,0x00),(0xC2,0x5E,0x00), g_websocket)

# ---------------------------------------------------------------------------
# Build the presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s

def textbox(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        for txt, size, color, bold, fname in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Consolas" if fname == "mono" else "Helvetica Neue"
    return tb

def kicker(s, txt):
    textbox(s, Inches(0.85), Inches(0.55), Inches(11), Inches(0.5),
            [[(txt.upper(), 15, ACCENT, True, None)]])

def title(s, txt, color=TEXT):
    textbox(s, Inches(0.85), Inches(0.95), Inches(11.6), Inches(1.1),
            [[(txt, 40, color, True, None)]])

def card(s, x, y, w, h, border=ACCENT):
    c = s.shapes.add_shape(5, x, y, w, h)  # rounded rect
    c.fill.solid(); c.fill.fore_color.rgb = CARD
    c.line.color.rgb = border; c.line.width = Pt(1.25)
    c.shadow.inherit = False
    try: c.adjustments[0] = 0.06
    except Exception: pass
    return c

def pic(s, path, x, y, w=None, h=None):
    return s.shapes.add_picture(path, x, y, width=w, height=h)

EMU = 914400
def fit(path, max_w_in, max_h_in):
    im = Image.open(path); iw, ih = im.size
    ar = iw / ih
    w, h = max_w_in, max_w_in / ar
    if h > max_h_in:
        h = max_h_in; w = max_h_in * ar
    return Inches(w), Inches(h)

# ---- Slide 1: Title ----
s = slide()
bar = s.shapes.add_shape(1, 0, 0, Inches(0.22), SH)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit=False
kicker(s, "EasyPilot · 4AHITM Drone Project")
textbox(s, Inches(0.85), Inches(2.2), Inches(11.5), Inches(2),
        [[("The EasyPilot ", 54, TEXT, True, None), ("iOS App", 54, ACCENT, True, None)]])
textbox(s, Inches(0.9), Inches(3.5), Inches(9.5), Inches(1.6),
        [[("A SwiftUI companion app that auto-discovers the drone on the network,", 20, MUTED, False, None)],
         [("streams live telemetry at 10 Hz, flies a 3D simulator, and sends", 20, MUTED, False, None)],
         [("flight commands over WebSocket.", 20, MUTED, False, None)]], sp_after=2)
# logo strip
xs = Inches(0.9); y = Inches(5.5)
for p, hh in [("swift.png",0.7),("swiftui.png",0.85),("scenekit.png",0.85),
              ("coremotion.png",0.85),("network.png",0.85),("esp.png",0.55)]:
    w,h = fit(os.path.join(LOGOS,p), 1.6, hh)
    img = pic(s, os.path.join(LOGOS,p), xs, y + (Inches(0.9)-h)//2, h=h)
    xs = Emu(int(xs) + int(img.width) + Inches(0.45))

# ---- Slide 2: Architecture ----
s = slide()
kicker(s, "System Overview")
title(s, "How the app talks to the drone")
# esp box
card(s, Inches(0.85), Inches(1.95), Inches(3.6), Inches(0.95), GREEN)
textbox(s, Inches(0.85), Inches(1.95), Inches(3.6), Inches(0.95),
        [[("ESP32-C3  ·  WiFi", 19, TEXT, True, None)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(5.0), Inches(2.0), Inches(7.6), Inches(0.9),
        [[("└─ UDP broadcast  ", 16, MUTED, False, None), ("\"EASYPILOT:<IP>\"", 16, GREEN, True, "mono"),
          ("  every 5 s · port 4242", 16, MUTED, False, None)],
         [("     → iOS auto-discovers the IP, no typing", 15, MUTED, False, None)]], sp_after=2)
textbox(s, Inches(5.0), Inches(3.05), Inches(7.6), Inches(0.9),
        [[("└─ WebSocket server · port 81", 16, ACCENT, True, None)],
         [("     ←→ commands out · telemetry JSON in @ 10 Hz", 15, MUTED, False, None)]], sp_after=2)
card(s, Inches(0.85), Inches(4.2), Inches(3.6), Inches(0.95), ACCENT)
textbox(s, Inches(0.85), Inches(4.2), Inches(3.6), Inches(0.95),
        [[("iPhone — EasyPilot App", 18, TEXT, True, None)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# stat cards
stats = [("3", "tabs: Dashboard · Simulator · Algorithms"),
         ("10 Hz", "telemetry & motion update rate"),
         ("~3,400", "lines of Swift across 18 files")]
cw = Inches(3.8); gap = Inches(0.3); x = Inches(0.85); y = Inches(5.55)
for big, lab in stats:
    card(s, x, y, cw, Inches(1.3))
    textbox(s, x, Inches(5.65), cw, Inches(0.7), [[(big, 32, ACCENT, True, None)]], align=PP_ALIGN.CENTER)
    textbox(s, x+Inches(0.2), Inches(6.35), cw-Inches(0.4), Inches(0.55), [[(lab, 13, MUTED, False, None)]], align=PP_ALIGN.CENTER)
    x = Emu(int(x)+int(cw)+int(gap))

# ---- Slide 3: Technologies (logos) ----
s = slide()
kicker(s, "Technologies Used")
title(s, "The iOS stack")
tech = [
    ("swift.png",      0.7,  "Swift 5", "Apple's modern, type-safe language for the whole app."),
    ("swiftui.png",    1.0,  "SwiftUI", "Declarative UI — TabView, glass cards, live bindings."),
    ("network.png",    1.0,  "Network.framework", "NWListener UDP discovery + WebSocket data link."),
    ("scenekit.png",   1.0,  "SceneKit", "Real-time 3D flight simulator with .usdz drone model."),
    ("coremotion.png", 1.0,  "CoreMotion", "Gyro / accelerometer read at 10 Hz for tilt control."),
    ("esp.png",        0.55, "Espressif ESP32", "On-drone firmware: WiFi, WebSocket server, beacon."),
]
cw, ch = Inches(3.85), Inches(2.1)
gx, gy = Inches(0.3), Inches(0.3)
x0, y0 = Inches(0.85), Inches(2.0)
for i,(img,hh,name,desc) in enumerate(tech):
    col, row = i % 3, i // 3
    x = Emu(int(x0)+col*(int(cw)+int(gx)))
    y = Emu(int(y0)+row*(int(ch)+int(gy)))
    card(s, x, y, cw, ch)
    w,h = fit(os.path.join(LOGOS,img), 1.2, hh)
    pic(s, os.path.join(LOGOS,img), x+Inches(0.28), y+Inches(0.32), h=h)
    textbox(s, x+Inches(1.7), y+Inches(0.3), cw-Inches(1.9), Inches(0.6),
            [[(name, 18, TEXT, True, None)]])
    textbox(s, x+Inches(0.3), y+Inches(1.15), cw-Inches(0.55), Inches(0.85),
            [[(desc, 13.5, MUTED, False, None)]], sp_after=0)

# ---- helper: code slide ----
def code_slide(kick, ttl, subtitle, code_lines):
    s = slide()
    kicker(s, kick); title(s, ttl)
    if subtitle:
        textbox(s, Inches(0.85), Inches(1.85), Inches(11.6), Inches(0.7),
                [[(subtitle, 16, MUTED, False, None)]])
    top = Inches(2.65) if subtitle else Inches(2.0)
    box = s.shapes.add_shape(5, Inches(0.85), top, Inches(11.6), Inches(4.3))
    box.fill.solid(); box.fill.fore_color.rgb = CODEBG
    box.line.color.rgb = ACCENT; box.line.width = Pt(1)
    box.shadow.inherit = False
    try: box.adjustments[0] = 0.03
    except Exception: pass
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.25); tf.margin_right = Inches(0.2)
    first = True
    for segs in code_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(2); p.line_spacing = 1.0
        for txt, col in segs:
            r = p.add_run(); r.text = txt
            r.font.name = "Consolas"; r.font.size = Pt(14.5); r.font.color.rgb = col
    return s

KW = RGBColor(0xFF,0x7A,0xB2); TY = RGBColor(0x7E,0xE0,0xC8)
ST = RGBColor(0xC3,0xE8,0x8D); CM = RGBColor(0x5D,0x6B,0x8F)
NU = RGBColor(0xF7,0xC8,0x73); FN = RGBColor(0x82,0xAA,0xFF); PL = TEXT

# ---- Slide 4: discovery code ----
code_slide("Code · Networking", "Zero-config discovery + WebSocket",
    "The app finds the drone with no manual IP entry — it listens for the UDP beacon, then upgrades to WebSocket.",
    [
     [("// Listen for the ESP32's UDP beacon on port 4242", CM)],
     [("private func ", KW),("receiveBeacon", FN),("(on conn: ", PL),("NWConnection", TY),(") {", PL)],
     [("    conn.", PL),("receiveMessage", FN),(" { [", PL),("weak self", KW),("] data, _, _, _ ", PL),("in", KW)],
     [("        if let", KW),(" data,", PL)],
     [("           let", KW),(" text = ", PL),("String", TY),("(data: data, encoding: .utf8),", PL)],
     [("           text.", PL),("hasPrefix", FN),("(", PL),("\"EASYPILOT:\"", ST),(") {", PL)],
     [("            let", KW),(" ip = ", PL),("String", TY),("(text.", PL),("dropFirst", FN),("(10))", PL)],
     [("            ", PL),("connectWebSocket", FN),("(to: ip)   ", PL),("// ws://<ip>:81", CM)],
     [("        }", PL)],
     [("    }", PL)],
     [("}", PL)],
    ])

# ---- Slide 5: telemetry code ----
code_slide("Code · Live Telemetry", "Decoding 10 Hz telemetry into the UI", None,
    [
     [("// Codable struct mirrors the ESP32 JSON exactly", CM)],
     [("struct ", KW),("DroneTelemetry", TY),(": ", PL),("Codable", TY),(" {", PL)],
     [("    let", KW),(" roll, pitch, yaw: ", PL),("Float", TY)],
     [("    let", KW),(" m1, m2, m3, m4: ", PL),("Int", TY),("?", PL)],
     [("    let", KW),(" voltage: ", PL),("Float", TY),("?, armed: ", PL),("Bool", TY),("?, mode: ", PL),("String", TY),("?", PL)],
     [("}", PL)],
     [("", PL)],
     [("private func ", KW),("decodeJSON", FN),("(_ text: ", PL),("String", TY),(") {", PL)],
     [("    let", KW),(" decoded = ", PL),("try", KW),(" ", PL),("JSONDecoder", TY),("().", PL),("decode", FN),("(", PL),("DroneTelemetry", TY),(".self, ...)", PL)],
     [("    ", PL),("DispatchQueue", TY),(".main.", PL),("async", FN),(" {", PL)],
     [("        self", KW),(".telemetry   = decoded   ", PL),("// @Published → SwiftUI redraws", CM)],
     [("        self", KW),(".isConnected = ", PL),("true", KW)],
     [("    }", PL),("\n}", PL)],
    ])

# ---- Slide 6: sensors + command ----
code_slide("Code · Sensors & Commands", "Phone tilt → drone command", None,
    [
     [("// CoreMotion device-motion updates @ 10 Hz", CM)],
     [("motionManager.", PL),("startDeviceMotionUpdates", FN),("(to: .main) { data, _ ", PL),("in", KW)],
     [("    let", KW),(" r = ", PL),("180", NU),(" / .pi", PL)],
     [("    self", KW),(".pitch = data.attitude.pitch * r", PL)],
     [("    self", KW),(".roll  = data.attitude.roll  * r", PL)],
     [("}", PL)],
     [("", PL)],
     [("// Type-safe JSON command builder (ControlProfile)", CM)],
     [("func ", KW),("startBalanceCommand", FN),("() -> ", PL),("String", TY),(" {", PL)],
     [("    ", PL),("\"\"\"{\"cmd\":\"START_BALANCE\",\"baseThrottle\":\\(baseThrottle),", ST)],
     [("       \"kPRoll\":\\(kPRoll),\"kPPitch\":\\(kPPitch)}\"\"\"", ST)],
     [("}", PL)],
    ])

# ---- Slide 7: Challenges ----
s = slide()
kicker(s, "Challenges")
title(s, "What was hard — and how we solved it")
ch_data = [
    ("Finding the drone", "The ESP32 gets a random DHCP IP. A UDP beacon broadcasts the IP every 5 s, so the app auto-connects — no typing.", GREEN),
    ("Detecting a dropped link", "WiFi can die silently. A 3 s timeout timer + 5 s WebSocket pings flip isConnected false the instant data stops.", ACCENT),
    ("Thread safety", "SceneKit physics runs on the render thread, SwiftUI only on main. Physics writes _vars, synced to @Published at 20 Hz.", ORANGE),
    ("Flight safety", "Accidental motor spin is dangerous. 1.5 s hold-to-arm, debounced safe-test, and a tilt-over-60° emergency stop.", RGBColor(0xFF,0x45,0x45)),
]
cw, chh = Inches(5.85), Inches(2.05)
x0, y0 = Inches(0.85), Inches(2.0)
for i,(t,desc,col) in enumerate(ch_data):
    cx = Emu(int(x0)+(i%2)*(int(cw)+int(Inches(0.3))))
    cy = Emu(int(y0)+(i//2)*(int(chh)+int(Inches(0.3))))
    card(s, cx, cy, cw, chh, col)
    textbox(s, cx+Inches(0.35), cy+Inches(0.3), cw-Inches(0.6), Inches(0.6),
            [[(t, 21, col, True, None)]])
    textbox(s, cx+Inches(0.35), cy+Inches(0.95), cw-Inches(0.7), Inches(1.0),
            [[(desc, 14.5, MUTED, False, None)]], sp_after=0)

# ---- Slide 8: threading code ----
code_slide("Code · The threading challenge", "Two threads, one drone",
    "Physics on SceneKit's render thread must never touch SwiftUI directly — so internal state is kept separate from published state.",
    [
     [("class ", KW),("DroneSimulator", TY),(": ", PL),("ObservableObject", TY),(" {", PL)],
     [("    // Main thread only — drives SwiftUI", CM)],
     [("    @", PL),("Published", TY),(" var", KW),(" simRoll = ", PL),("0.0", NU)],
     [("    @", PL),("Published", TY),(" var", KW),(" isCrashed = ", PL),("false", KW)],
     [("", PL)],
     [("    // Render thread only — raw physics integration", CM)],
     [("    var", KW),(" _roll = ", PL),("0.0", NU)],
     [("    var", KW),(" _vel: ", PL),("SIMD3", TY),("<", PL),("Float", TY),("> = .zero", PL)],
     [("", PL)],
     [("    func ", KW),("tick", FN),("(dt: ", PL),("Double", TY),(") { ", PL),("/* render thread: physics */", CM),(" }", PL)],
     [("    func ", KW),("syncPublished", FN),("() { ", PL),("/* main: _vars → @Published */", CM),(" }", PL)],
     [("}", PL)],
    ])

# ---- Slide 9: Summary ----
s = slide()
kicker(s, "Wrap-up")
title(s, "Where we are & what's next")
card(s, Inches(0.85), Inches(2.0), Inches(5.85), Inches(3.5), GREEN)
textbox(s, Inches(1.2), Inches(2.25), Inches(5), Inches(0.6), [[("✓  Done", 24, GREEN, True, None)]])
done = ["Auto-discovery + WebSocket telemetry","Live Dashboard with 3D drone view",
        "Full 3D flight Simulator (Rate / Balance / PosHold)","Mode 2 virtual joysticks",
        "Tunable, savable control profiles"]
textbox(s, Inches(1.25), Inches(3.1), Inches(5.2), Inches(2.3),
        [[("•  "+d, 16, TEXT, False, None)] for d in done], sp_after=8)
card(s, Inches(7.0), Inches(2.0), Inches(5.45), Inches(3.5), ORANGE)
textbox(s, Inches(7.35), Inches(2.25), Inches(5), Inches(0.6), [[("→  Next (Sprint 6)", 24, ORANGE, True, None)]])
nxt = ["Live mode — fly the real drone from the joysticks","Real MSP telemetry from the Betaflight FC",
       "Landscape UI polish & touch-aware stabilization"]
textbox(s, Inches(7.4), Inches(3.1), Inches(4.8), Inches(2.3),
        [[("•  "+d, 16, TEXT, False, None)] for d in nxt], sp_after=10)
textbox(s, Inches(0.85), Inches(6.4), Inches(11.6), Inches(0.6),
        [[("EasyPilot · 4AHITM HTL · iOS subsystem overview", 15, MUTED, False, None)]], align=PP_ALIGN.CENTER)

out = os.path.join(HERE, "EasyPilot-iOS-Overview.pptx")
prs.save(out)
print("Saved:", out, "—", len(prs.slides._sldIdLst), "slides")
