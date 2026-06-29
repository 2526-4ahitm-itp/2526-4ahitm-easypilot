#!/usr/bin/env python3
"""Builds the 'EasyPilot iOS' PowerPoint decks — own modern style.

Produces two files from the same slide definitions:
  EasyPilot-iOS.pptx           4 slides (Title, Was ist, Architektur, Live Demo)
  EasyPilot-iOS-Extended.pptx  5 slides (adds a "Das Projekt EasyPilot" slide)
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE   = os.path.dirname(os.path.abspath(__file__))
LOGOS  = os.path.join(HERE, "logos")
OUT     = os.path.join(HERE, "EasyPilot-iOS.pptx")
OUT_EXT = os.path.join(HERE, "EasyPilot-iOS-Extended.pptx")

NAVY   = RGBColor(0x22, 0x2B, 0x3C)
TEAL   = RGBColor(0x16, 0x21, 0x33)
ACCENT = RGBColor(0x2F, 0x80, 0xED)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x2B, 0x33, 0x40)
MUTED  = RGBColor(0x6B, 0x74, 0x82)
LIGHT  = RGBColor(0xC4, 0xCE, 0xDE)
CARD   = RGBColor(0xF4, 0xF7, 0xFB)
BORDER = RGBColor(0xE2, 0xE7, 0xEE)

HEAD = "Segoe UI Semibold"
SANS = "Segoe UI"
RR   = MSO_SHAPE.ROUNDED_RECTANGLE
SW   = Inches(13.333)
SH   = Inches(7.5)

prs = None  # current presentation (set by build())


def new_prs():
    p = Presentation()
    p.slide_width  = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def slide(bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def textbox(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            sp_after=8, line=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        if line:
            p.line_spacing = line
        for txt, size, color, bold, fname in para:
            run = p.add_run(); run.text = txt
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = fname
    return tb


def rect(s, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, border=None, bw=1.0, rad=None):
    c = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        c.fill.background()
    else:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if border is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = border; c.line.width = Pt(bw)
    c.shadow.inherit = False
    if rad is not None:
        try: c.adjustments[0] = rad
        except Exception: pass
    return c


def connector(s, x1, y1, x2, y2, col=NAVY, pt=2.25, head=True):
    a = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    a.line.color.rgb = col; a.line.width = Pt(pt); a.shadow.inherit = False
    if head:
        ln = a.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return a


def fit(path, maxw_in, maxh_in):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = maxw_in, maxw_in / ar
    if h > maxh_in:
        h, w = maxh_in, maxh_in * ar
    return w, h


def icon(s, path, cx_in, top_in, box_w, box_h):
    w, h = fit(os.path.join(LOGOS, path), box_w, box_h)
    x = Inches(cx_in - w / 2)
    y = Inches(top_in + (box_h - h) / 2)
    s.shapes.add_picture(os.path.join(LOGOS, path), x, y, width=Inches(w), height=Inches(h))


def header_band(s, title, num):
    rect(s, 0, 0, SW, Inches(1.35), NAVY)
    rect(s, Inches(0.85), Inches(0.5), Inches(0.34), Inches(0.34), ACCENT)
    textbox(s, Inches(1.4), 0, Inches(10), Inches(1.35),
            [[(title, 30, WHITE, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(12.0), 0, Inches(0.9), Inches(1.35),
            [[(num, 16, ACCENT, True, HEAD)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def statement_block(s, statement, paras):
    """Big key statement + calm tick paragraphs (left column, right left empty)."""
    textbox(s, Inches(0.95), Inches(1.85), Inches(8.2), Inches(1.4),
            [[(statement, 27, TEAL, True, HEAD)]], line=1.08)
    rect(s, Inches(0.97), Inches(3.12), Inches(1.6), Pt(4), ACCENT)
    for y, barh, body in paras:
        rect(s, Inches(0.95), Inches(y + 0.03), Inches(0.06), Inches(barh), ACCENT)
        textbox(s, Inches(1.3), Inches(y), Inches(7.4), Inches(1.5),
                [[(body, 16.5, INK, False, SANS)]], line=1.16)


# --------------------------------------------------------------------- slides
def title_slide():
    s = slide(NAVY)
    rect(s, 0, 0, Inches(0.16), SH, ACCENT)
    textbox(s, Inches(0.95), Inches(1.5), Inches(11), Inches(0.5),
            [[("EASYPILOT  ·  4AHITM DROHNENPROJEKT", 14, ACCENT, True, HEAD)]])
    textbox(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(1.5),
            [[("EasyPilot ", 62, WHITE, True, HEAD), ("iOS", 62, ACCENT, True, HEAD)]])
    rect(s, Inches(0.95), Inches(3.62), Inches(1.5), Pt(4), ACCENT)
    textbox(s, Inches(0.95), Inches(3.95), Inches(10.8), Inches(1.2),
            [[("SwiftUI-Companion-App für die EasyPilot-Drohne – findet die Drohne", 20, LIGHT, False, SANS)],
             [("selbst im WLAN, zeigt Live-Telemetrie und einen 3D-Flugsimulator.", 20, LIGHT, False, SANS)]],
            sp_after=3)
    textbox(s, Inches(0.95), Inches(6.1), Inches(8), Inches(0.9),
            [[("Simon Eder", 19, WHITE, True, HEAD)],
             [("Creator", 14, LIGHT, False, SANS)]], sp_after=2)


def purpose_slide(num):
    s = slide(WHITE)
    header_band(s, "Das Projekt EasyPilot", num)
    statement_block(
        s,
        "Eine gewöhnliche Drohne wird zur vernetzten Telemetrie- und Steuerplattform.",
        [
            (3.55, 0.62, "Ein ESP32-C3 an Bord macht die Drohne über WLAN ansprechbar – "
                         "er sendet Telemetrie und nimmt Flugbefehle entgegen."),
            (4.62, 0.62, "Mehrere Clients sind live dabei: eine iOS-App und ein 3D-Web-Viewer "
                         "zeigen Lage und Daten der Drohne in Echtzeit."),
            (5.69, 0.62, "Ein Schulprojekt der 4AHITM (HTL Leonding) – es verbindet "
                         "Embedded-Firmware, Echtzeit-Netzwerk sowie App- und Web-Entwicklung."),
        ],
    )


def wasist_slide(num, algo=False):
    s = slide(WHITE)
    header_band(s, "Was ist EasyPilot iOS", num)
    if algo:
        paras = [
            (3.55, 0.62, "Verbindet sich von selbst über das WLAN mit der Drohne – ohne Eingabe einer "
                         "IP-Adresse."),
            (4.70, 0.62, "Zeigt die Telemetrie zehnmal pro Sekunde in Echtzeit an und bringt einen "
                         "eigenen 3D-Flugsimulator direkt aufs Handy."),
            (5.85, 0.62, "Stabilisierungs- und Hilfsalgorithmen lassen sich durch die App steuern."),
        ]
    else:
        paras = [
            (3.70, 0.62, "Verbindet sich von selbst über das WLAN mit der Drohne – ohne Eingabe einer "
                         "IP-Adresse."),
            (5.00, 0.62, "Zeigt die Telemetrie zehnmal pro Sekunde in Echtzeit an und bringt einen "
                         "eigenen 3D-Flugsimulator direkt aufs Handy."),
        ]
    statement_block(
        s, "Das iPhone wird zur Anzeige- und Steuerzentrale für unsere Drohne.", paras)


def arch_slide(num):
    s = slide(WHITE)
    header_band(s, "Architektur & Tech-Stack", num)

    # LEFT — the drone
    rect(s, Inches(0.85), Inches(2.45), Inches(2.95), Inches(3.0), CARD, RR, border=BORDER, bw=1, rad=0.05)
    textbox(s, Inches(0.85), Inches(2.62), Inches(2.95), Inches(0.4),
            [[("DROHNE", 13, MUTED, True, HEAD)]], align=PP_ALIGN.CENTER)
    icon(s, "chip.png", 0.85 + 1.475, 3.05, 1.25, 0.9)
    textbox(s, Inches(0.85), Inches(4.1), Inches(2.95), Inches(0.45),
            [[("ESP32-C3", 18, TEAL, True, HEAD)]], align=PP_ALIGN.CENTER)
    textbox(s, Inches(0.85), Inches(4.55), Inches(2.95), Inches(0.8),
            [[("WebSocket-Server", 12.5, MUTED, False, SANS)],
             [("UDP-Beacon (Discovery)", 12.5, MUTED, False, SANS)]],
            align=PP_ALIGN.CENTER, sp_after=1)

    # ARROW drone -> app
    connector(s, Inches(3.8), Inches(3.62), Inches(5.35), Inches(3.62))
    textbox(s, Inches(3.7), Inches(2.95), Inches(1.75), Inches(0.6),
            [[("WLAN · WebSocket", 11.5, MUTED, True, HEAD)]], align=PP_ALIGN.CENTER)
    textbox(s, Inches(3.7), Inches(3.78), Inches(1.75), Inches(0.5),
            [[("Telemetrie · 10 Hz", 11.5, MUTED, False, SANS)]], align=PP_ALIGN.CENTER)

    # RIGHT — the iPhone app
    APP_X, APP_W = 5.35, 7.05
    rect(s, Inches(APP_X), Inches(1.95), Inches(APP_W), Inches(4.95), WHITE, RR, border=ACCENT, bw=1.75, rad=0.04)
    textbox(s, Inches(APP_X + 0.3), Inches(2.12), Inches(4.2), Inches(0.45),
            [[("iPhone-App", 18, TEAL, True, HEAD)]])
    textbox(s, Inches(APP_X + 0.3), Inches(2.58), Inches(4.2), Inches(0.4),
            [[("gebaut mit Swift & SwiftUI", 12.5, MUTED, False, SANS)]])
    icon(s, "swift.png",   APP_X + APP_W - 1.65, 2.18, 1.25, 0.55)
    icon(s, "swiftui.png", APP_X + APP_W - 0.55, 2.12, 0.62, 0.62)

    comps = [
        ("wifi.png", "Network.framework", "Verbindung & Auto-Discovery zur Drohne"),
        ("cube3d.png", "SceneKit",         "Rendert den 3D-Flugsimulator"),
        ("gyro.png", "CoreMotion",         "Liest das Gyroskop des Handys"),
    ]
    cx, cw, chh = APP_X + 0.3, APP_W - 0.6, 1.05
    cy = 3.15
    for img, name, sub in comps:
        rect(s, Inches(cx), Inches(cy), Inches(cw), Inches(chh), CARD, RR, border=BORDER, bw=1, rad=0.08)
        icon(s, img, cx + 0.7, cy + 0.12, 0.95, chh - 0.24)
        textbox(s, Inches(cx + 1.4), Inches(cy + 0.18), Inches(cw - 1.6), Inches(0.45),
                [[(name, 17, TEAL, True, HEAD)]])
        textbox(s, Inches(cx + 1.4), Inches(cy + 0.6), Inches(cw - 1.6), Inches(0.4),
                [[(sub, 12.5, MUTED, False, SANS)]])
        cy += chh + 0.13


def demo_slide(num):
    s = slide(NAVY)
    rect(s, 0, 0, Inches(0.16), SH, ACCENT)
    textbox(s, Inches(0.95), Inches(2.7), Inches(8), Inches(0.5),
            [[(f"{num}  ·  DEMO", 14, ACCENT, True, HEAD)]])
    textbox(s, Inches(0.9), Inches(3.2), Inches(10), Inches(1.4),
            [[("Live Demo", 66, WHITE, True, HEAD)]])
    rect(s, Inches(0.95), Inches(4.75), Inches(1.6), Pt(4), ACCENT)


def build(out, extended):
    global prs
    prs = new_prs()
    title_slide()
    if extended:
        purpose_slide("02")
        wasist_slide("03"); arch_slide("04"); demo_slide("05")
    else:
        wasist_slide("02", algo=True); arch_slide("03"); demo_slide("04")
    prs.save(out)
    print("saved", out)


build(OUT, False)
build(OUT_EXT, True)
